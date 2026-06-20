# -*- coding: utf-8 -*-
"""生成课程汇报 PPT：数据 → 模型 → 结果（智农项目）"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.join(ROOT, "docs", "智农项目_课程汇报.pptx")

# Forest & Moss — 与项目农业主题一致
C_FOREST = RGBColor(0x2C, 0x5F, 0x2D)
C_MOSS = RGBColor(0x97, 0xBC, 0x62)
C_CREAM = RGBColor(0xF5, 0xF5, 0xF0)
C_DARK = RGBColor(0x1A, 0x2E, 0x1A)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED = RGBColor(0x4A, 0x5D, 0x4A)
C_ACCENT = RGBColor(0xB8, 0x50, 0x42)


def set_slide_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_FOREST
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(9), Inches(0.55))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Microsoft YaHei"
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.72), Inches(9), Inches(0.35))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = C_MOSS
        sp.font.name = "Microsoft YaHei"


def add_bullets(slide, items, left=0.55, top=1.35, width=8.9, height=4.0, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = C_DARK
        p.space_after = Pt(8)


def add_stat_block(slide, x, y, w, h, num, label, fill=C_MOSS):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = C_FOREST
    tb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.15), Inches(w - 0.2), Inches(h - 0.2))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.text = str(num)
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = C_FOREST
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "Microsoft YaHei"
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = C_DARK
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Microsoft YaHei"


def build():
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    blank = prs.slide_layouts[6]

    # --- 1 封面 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_FOREST)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.85), Inches(10), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_MOSS
    accent.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = "智农 · 农作物病虫害 AI 智能识别与预警"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.font.name = "Microsoft YaHei"
    t2 = s.shapes.add_textbox(Inches(0.6), Inches(2.75), Inches(8.8), Inches(0.8))
    p2 = t2.text_frame.paragraphs[0]
    p2.text = "课程汇报：从数据到模型到结果的机器学习实践"
    p2.font.size = Pt(18)
    p2.font.color.rgb = C_MOSS
    p2.font.name = "Microsoft YaHei"
    t3 = s.shapes.add_textbox(Inches(0.6), Inches(4.0), Inches(8.8), Inches(0.5))
    p3 = t3.text_frame.paragraphs[0]
    p3.text = "多任务深度学习 · PlantVillage · Flask Web 部署"
    p3.font.size = Pt(14)
    p3.font.color.rgb = C_CREAM
    p3.font.name = "Microsoft YaHei"

    # --- 2 汇报路线 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "汇报思路", "符合课程要求：数据 → 模型 → 结果")
    add_bullets(s, [
        "① 数据：来源、收集、清洗、增强、训练/验证划分",
        "② 模型：多任务学习思想、网络结构、损失与训练流程",
        "③ 调优：迁移学习、早停、学习率与类别不平衡处理",
        "④ 结果：Accuracy / F1 / Macro-F1 等定量指标与系统演示",
        "⑤ 拓展：Grad-CAM 可解释性、MC Dropout 不确定性、防治知识库",
    ])

    # --- 3 背景 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "1. 项目背景与问题定义")
    add_bullets(s, [
        "农业场景：农户上传叶片照片，需要快速判断「是什么病、有多严重、风险多高」。",
        "传统方式依赖专家经验，响应慢、成本高，难以覆盖 60+ 常见作物病害。",
        "机器学习任务：图像分类 + 多任务联合（病害类别 + 严重程度 + 风险评分）。",
        "项目目标：端到端系统 — 数据训练 → 模型推理 → Web 交互 → 可执行防治方案。",
    ])

    # --- 4 数据介绍 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "2. 数据介绍", "公开数据集与项目自建流程")
    add_bullets(s, [
        "主数据集：PlantVillage — 按「作物_病害」文件夹组织，约 61 类叶片图像。",
        "演示/快速验证：scripts/generate_mock_dataset.py 生成小规模多类样本。",
        "Kaggle 流程：prepare_kaggle_dataset.py 切分为 train/val，便于客户演示。",
        "输入规格：RGB 图像，推理时统一 resize 至 128×128（与训练一致）。",
        "标签体系：病害 class_id；严重程度由类别名/规则映射为 健康 / 一般 / 严重。",
    ], top=1.25)
    add_stat_block(s, 0.55, 4.35, 2.0, 0.95, "61", "模型类别（PlantVillage）")
    add_stat_block(s, 2.75, 4.35, 2.0, 0.95, "60+", "防治知识库条目")
    add_stat_block(s, 4.95, 4.35, 2.0, 0.95, "128²", "输入分辨率")
    add_stat_block(s, 7.15, 4.35, 2.3, 0.95, "train/val", "标准划分")

    # --- 5 数据处理 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "3. 数据收集、处理与划分")
    add_bullets(s, [
        "收集：ImageFolder / TXT 标注双模式；build_image_index 建立文件名索引。",
        "清洗：文件签名检测、拒绝非图片；损坏样本返回占位图避免训练中断。",
        "增强（训练集）：RandomResizedCrop、水平翻转、ColorJitter、Normalize（ImageNet 统计量）。",
        "划分：公开集按 train/val 目录；可 sample_ratio 子采样做快速实验。",
        "类别映射：disease_mapping 将文件夹名映射为整数标签，保存至 best_multitask_model.json。",
    ])

    # --- 6 方法概述 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "4. 机器学习方法概述", "为什么选择多任务学习？")
    add_bullets(s, [
        "单任务：仅预测病害 OR 仅预测严重程度 — 两个独立模型，部署与维护成本高。",
        "多任务学习（MTL）：共享 MobileNetV2 骨干，联合优化病害头 + 严重程度头。",
        "优势：相关任务共享表征，通常提升小样本类泛化；一次前向得到多维度决策。",
        "项目还实现 SingleTaskDiseaseModel / SingleTaskSeverityModel 用于对比实验。",
        "推理阶段：softmax 得类别与严重度 → 规则/模型融合得 risk_score（0–100）。",
    ])

    # --- 7 模型结构 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "5. 模型结构设计")
    add_bullets(s, [
        "Backbone：MobileNetV2（ImageNet 预训练），轻量、适合 CPU/边缘部署。",
        "迁移学习：冻结 features 除最后 3 层，减少过拟合与训练时间。",
        "CrossTaskAttention：对 1280 维共享特征做病害/严重度双路注意力，输出 256 维。",
        "双头输出：disease_head → num_diseases；severity_head → 3 类。",
        "可解释：保留 feature_extractor 输出，支持 Grad-CAM 热力图（train_severity_gradcam.py）。",
    ], top=1.2, height=3.2)
    flow = s.shapes.add_textbox(Inches(0.55), Inches(4.35), Inches(8.9), Inches(0.9))
    fp = flow.text_frame.paragraphs[0]
    fp.text = "数据流：图像 → MobileNetV2 → GAP → CrossTaskAttention → [病害 logits | 严重度 logits]"
    fp.font.size = Pt(13)
    fp.font.italic = True
    fp.font.color.rgb = C_ACCENT
    fp.font.name = "Microsoft YaHei"

    # --- 8 训练 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "6. 模型训练与损失设计")
    add_bullets(s, [
        "优化器：Adam；损失：CrossEntropyLoss（病害 + 严重度加权求和）。",
        "Trainer：TaskTrainer 统一 multitask / 单任务训练循环，保证对比公平。",
        "命令示例：python train_multitask_model.py --dataset-mode plantvillage --epochs 15 --patience 5",
        "模拟数据快速验证：--data-dir data/mock_problem_b --epochs 5 --patience 2",
        "产出：best_multitask_model.pth + JSON（class_names、输入尺寸、指标记录）。",
    ])

    # --- 9 调优 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "7. 调参与实验策略")
    add_bullets(s, [
        "早停（Early Stopping）：验证集指标不再提升时停止，patience 可配置。",
        "Dropout 0.3/0.2：缓解过拟合，尤其类内样本不均衡时。",
        "MC Dropout（推理可选）：多次前向采样估计预测不确定性，低置信度提示人工复核。",
        "类别不平衡：Macro-F1 关注少数类；可结合加权 CE 或重采样（实验扩展）。",
        "超参：batch size、学习率、sample_ratio 用于课堂/演示的快速迭代。",
    ])

    # --- 10 评估指标 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "8. 定量精度分析", "sklearn.metrics + 验证集评估")
    add_bullets(s, [
        "病害分类：Accuracy（整体正确率）、Weighted F1（按样本量加权）。",
        "严重程度：Macro-F1（三类等权，关注「严重」漏检成本）。",
        "训练脚本在验证阶段输出 disease_acc、f1、severity macro-f1，并写入日志/JSON。",
        "业务指标：风险分档（低/中/高）、与防治 urgency 联动，形成可执行建议。",
        "说明：最终数值取决于数据规模与 epoch；完整 PlantVillage 训练通常 Acc 可达较高水平（需本地复现实验填表）。",
    ], top=1.2, height=2.8)
    tbl = s.shapes.add_table(4, 4, Inches(0.55), Inches(3.95), Inches(8.9), Inches(1.35)).table
    headers = ["指标", "含义", "适用任务", "备注"]
    rows = [
        ["Accuracy", "预测正确比例", "病害分类", "直观、易汇报"],
        ["Weighted F1", "精确率召回率调和平均", "病害分类", "适合不均衡"],
        ["Macro-F1", "各类 F1 算术平均", "严重度 3 类", "重视少数类"],
    ]
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            tbl.cell(r, c).text = val
            tbl.cell(r, c).text_frame.paragraphs[0].font.size = Pt(10)

    # --- 11 结果与系统 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "9. 实验结果与系统落地")
    add_bullets(s, [
        "离线：混淆矩阵、训练曲线（train_multitask_model.py 可视化模块）。",
        "在线：Flask app.py — 上传/URL/批量预测、JSON·PDF 报告、ECharts 风险趋势。",
        "知识库：60+ 病害防治方案（症状、用药、复查时间线），与识别结果自动匹配。",
        "演示模式：无模型亦可展示 demo_case + 7 条历史报告，适合答辩与客户展示。",
        "扩展：disease_multitask.pth 切换环境变量 MODEL_PATH 即可对比不同 checkpoint。",
    ])

    # --- 12 流程图式总结 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "10. 端到端流程回顾")
    steps = [
        ("数据", "PlantVillage / Kaggle / Mock\n增强 + train/val"),
        ("模型", "MobileNetV2 + MTL\nCrossTaskAttention"),
        ("训练", "CE Loss + Adam\n早停 + 迁移学习"),
        ("评估", "Acc / F1 / Macro-F1\n混淆矩阵"),
        ("部署", "Flask + 知识库\nGrad-CAM / MC Dropout"),
    ]
    for i, (title, body) in enumerate(steps):
        x = 0.4 + i * 1.92
        fill = C_MOSS if i % 2 == 0 else RGBColor(0xE7, 0xE8, 0xD1)
        add_stat_block(s, x, 1.45, 1.75, 1.35, str(i + 1), title, fill=fill)
        bx = s.shapes.add_textbox(Inches(x + 0.08), Inches(2.95), Inches(1.6), Inches(1.5))
        bp = bx.text_frame.paragraphs[0]
        bp.text = body
        bp.font.size = Pt(10)
        bp.font.name = "Microsoft YaHei"
        bp.font.color.rgb = C_DARK
    add_bullets(s, [
        "汇报结论：本项目完整走通「数据准备 → 模型设计训练 → 定量评估 → 工程化演示」机器学习闭环。",
    ], top=4.55, height=0.8, size=14)

    # --- 13 展望 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_CREAM)
    add_title_bar(s, "11. 总结与展望")
    add_bullets(s, [
        "总结：多任务 CNN 在农业病害场景可同时输出类别与严重度，并结合规则生成风险与防治方案。",
        "不足：田间光照/遮挡与实验室数据集存在域差异；需持续收集本地数据微调。",
        "展望：更强骨干（EfficientNet/ViT）、半监督学习、移动端 ONNX 部署、与农技平台对接。",
        "课程关联：监督学习、卷积神经网络、迁移学习、多任务学习、模型评估与工程部署。",
    ])

    # --- 14 致谢 ---
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, C_FOREST)
    t = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(1.0))
    p = t.text_frame.paragraphs[0]
    p.text = "谢谢老师与同学"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"
    t2 = s.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(8.8), Inches(0.6))
    p2 = t2.text_frame.paragraphs[0]
    p2.text = "Q & A  ·  演示地址：http://localhost:7860"
    p2.font.size = Pt(16)
    p2.font.color.rgb = C_MOSS
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Microsoft YaHei"

    prs.save(OUT)
    print("Saved:", OUT)


if __name__ == "__main__":
    build()